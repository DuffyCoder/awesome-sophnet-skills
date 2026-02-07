#!/usr/bin/env python3
"""Create DOCX or PPTX files from inline text or markdown."""

from __future__ import annotations

import argparse
import re
import sys
from datetime import datetime
from pathlib import Path
from typing import List, Tuple


def fail(message: str) -> None:
    print(f"Error: {message}", file=sys.stderr)
    raise SystemExit(1)


def skill_root() -> Path:
    return Path(__file__).resolve().parent.parent


def uv_sync_hint() -> str:
    return f"uv sync --project {skill_root()}"


def sanitize_filename(raw: str, fallback: str) -> str:
    cleaned = re.sub(r"[^A-Za-z0-9_-]+", "_", raw).strip("_")
    return cleaned or fallback


def output_path_for(title: str, file_type: str, explicit_output: str | None) -> Path:
    if explicit_output:
        return Path(explicit_output).expanduser().resolve()
    stamp = datetime.now().strftime("%Y%m%d_%H%M%S")
    base = sanitize_filename(title, "document")
    return Path(f"{base}_{stamp}.{file_type}").resolve()


def load_content(inline_content: str | None, markdown_path: str | None) -> str:
    if markdown_path:
        path = Path(markdown_path).expanduser().resolve()
        if not path.is_file():
            fail(f"markdown file does not exist: {path}")
        return path.read_text(encoding="utf-8")
    return inline_content or ""


def parse_markdown_sections(text: str) -> List[Tuple[str, List[str]]]:
    sections: List[Tuple[str, List[str]]] = []
    current_title = "Overview"
    bullets: List[str] = []

    def flush() -> None:
        nonlocal bullets, current_title
        if bullets:
            sections.append((current_title, bullets))
            bullets = []

    for raw in text.splitlines():
        line = raw.strip()
        if not line:
            continue
        if line.startswith("#"):
            flush()
            current_title = line.lstrip("#").strip() or "Section"
            continue
        if line.startswith("- ") or line.startswith("* "):
            bullets.append(line[2:].strip())
            continue
        bullets.append(line)

    flush()
    return sections


def create_docx(title: str, content: str, author: str, destination: Path) -> None:
    try:
        from docx import Document
        from docx.enum.text import WD_ALIGN_PARAGRAPH
        from docx.shared import Pt
    except ImportError:
        fail(
            "python-docx is missing in the skill environment. "
            f"Run: {uv_sync_hint()}"
        )

    document = Document()
    style = document.styles["Normal"]
    style.font.name = "Arial"
    style.font.size = Pt(11)

    title_paragraph = document.add_heading(title, level=0)
    title_paragraph.alignment = WD_ALIGN_PARAGRAPH.CENTER

    document.add_paragraph(f"Author: {author}")
    document.add_paragraph(f"Created: {datetime.now().strftime('%Y-%m-%d %H:%M:%S')}")

    for raw in content.splitlines():
        line = raw.strip()
        if not line:
            continue
        if line.startswith("### "):
            document.add_heading(line[4:].strip(), level=3)
        elif line.startswith("## "):
            document.add_heading(line[3:].strip(), level=2)
        elif line.startswith("# "):
            document.add_heading(line[2:].strip(), level=1)
        elif line.startswith("- ") or line.startswith("* "):
            document.add_paragraph(line[2:].strip(), style="List Bullet")
        else:
            document.add_paragraph(line)

    destination.parent.mkdir(parents=True, exist_ok=True)
    document.save(str(destination))


def create_pptx(title: str, content: str, slides: int, destination: Path) -> None:
    try:
        from pptx import Presentation
    except ImportError:
        fail(
            "python-pptx is missing in the skill environment. "
            f"Run: {uv_sync_hint()}"
        )

    if slides < 1:
        fail("--slides must be >= 1")

    presentation = Presentation()

    title_slide = presentation.slides.add_slide(presentation.slide_layouts[0])
    title_slide.shapes.title.text = title
    if len(title_slide.placeholders) > 1:
        title_slide.placeholders[1].text = f"Created: {datetime.now().strftime('%Y-%m-%d')}"

    sections = parse_markdown_sections(content)

    if not sections:
        sections = [("Agenda", ["Introduction", "Key points", "Summary"])]

    max_content_slides = slides - 1
    if max_content_slides < len(sections):
        sections = sections[:max_content_slides]

    while len(sections) < max_content_slides:
        sections.append((f"Slide {len(sections) + 1}", ["Add content"]))

    for section_title, bullets in sections:
        slide = presentation.slides.add_slide(presentation.slide_layouts[1])
        slide.shapes.title.text = section_title
        body = slide.placeholders[1].text_frame
        body.clear()
        for index, bullet in enumerate(bullets):
            paragraph = body.paragraphs[0] if index == 0 else body.add_paragraph()
            paragraph.text = bullet
            paragraph.level = 0

    destination.parent.mkdir(parents=True, exist_ok=True)
    presentation.save(str(destination))


def parse_args() -> argparse.Namespace:
    parser = argparse.ArgumentParser(description="Create DOCX/PPTX files from text or markdown")
    parser.add_argument("type", choices=["docx", "pptx"], help="Document type")
    parser.add_argument("--title", required=True, help="Document title")
    parser.add_argument("--content", help="Inline markdown/plain content")
    parser.add_argument("--markdown", help="Path to markdown file")
    parser.add_argument("--author", default="Codex Assistant", help="Author label for DOCX")
    parser.add_argument("--slides", type=int, default=6, help="Slide count for PPTX")
    parser.add_argument("--output", help="Explicit output file path")
    return parser.parse_args()


def main() -> None:
    args = parse_args()
    content = load_content(args.content, args.markdown)
    destination = output_path_for(args.title, args.type, args.output)

    if args.type == "docx":
        create_docx(args.title, content, args.author, destination)
    else:
        create_pptx(args.title, content, args.slides, destination)

    print(f"FILE_PATH={destination}")
    print(f"FILE_TYPE={args.type}")
    print(f"FILE_NAME={destination.name}")


if __name__ == "__main__":
    main()
