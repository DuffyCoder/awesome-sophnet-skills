#!/usr/bin/env python3
"""Create an AI Introduction PowerPoint presentation using python-pptx"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

# Create presentation
prs = Presentation()
prs.slide_width = Inches(10)
prs.slide_height = Inches(5.625)

# Define blank slide layout (index 6 is typically blank)
blank_layout = prs.slide_layouts[6]

# Color palette - Teal Trust theme
class Colors:
    primary = RGBColor(2, 128, 144)      # Teal
    secondary = RGBColor(0, 168, 150)    # Seafoam
    accent = RGBColor(2, 195, 154)       # Mint
    dark = RGBColor(15, 23, 42)          # Dark slate
    light = RGBColor(248, 250, 252)      # Off-white
    white = RGBColor(255, 255, 255)
    gray = RGBColor(100, 116, 139)

# Helper functions
def add_title_slide(prs, blank_layout):
    """Slide 1: Title Slide"""
    slide = prs.slides.add_slide(blank_layout)

    # Dark background
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = Colors.dark

    # Decorative accent circle (using shape)
    shape = slide.shapes.add_shape(
        9,  # Oval
        Inches(7.5), Inches(-1.5), Inches(6), Inches(6)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = Colors.primary
    shape.line.fill.background()

    # Subtitle
    textbox = slide.shapes.add_textbox(Inches(0.5), Inches(1.8), Inches(9), Inches(0.6))
    text_frame = textbox.text_frame
    text_frame.text = "Introduction to"
    p = text_frame.paragraphs[0]
    p.font.size = Pt(28)
    p.font.color.rgb = Colors.secondary
    p.font.bold = True
    p.alignment = PP_ALIGN.LEFT

    # Main title
    textbox = slide.shapes.add_textbox(Inches(0.5), Inches(2.5), Inches(9), Inches(1.2))
    text_frame = textbox.text_frame
    text_frame.text = "Artificial Intelligence"
    p = text_frame.paragraphs[0]
    p.font.size = Pt(54)
    p.font.color.rgb = Colors.white
    p.font.bold = True
    p.alignment = PP_ALIGN.LEFT

    # Subtitle line
    textbox = slide.shapes.add_textbox(Inches(0.5), Inches(3.9), Inches(9), Inches(0.5))
    text_frame = textbox.text_frame
    text_frame.text = "Understanding the Technology Shaping Our Future"
    p = text_frame.paragraphs[0]
    p.font.size = Pt(18)
    p.font.color.rgb = Colors.gray
    p.alignment = PP_ALIGN.LEFT

    # Accent bar
    shape = slide.shapes.add_shape(
        1,  # Rectangle
        Inches(0.5), Inches(4.6), Inches(2), Inches(0.08)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = Colors.accent
    shape.line.fill.background()

def add_content_slide(prs, title, content_items, blank_layout, is_two_col=False):
    """Generic content slide with bullet points"""
    slide = prs.slides.add_slide(blank_layout)

    # Light background
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = Colors.light

    # Accent bar on left
    shape = slide.shapes.add_shape(
        1,  # Rectangle
        Inches(0), Inches(0), Inches(0.15), Inches(5.625)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = Colors.primary
    shape.line.fill.background()

    # Title
    textbox = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.7))
    text_frame = textbox.text_frame
    text_frame.text = title
    p = text_frame.paragraphs[0]
    p.font.size = Pt(44)
    p.font.color.rgb = Colors.dark
    p.font.bold = True
    p.alignment = PP_ALIGN.LEFT

    if is_two_col:
        # Two column layout
        left_x = Inches(0.5)
        right_x = Inches(5)
        width = Inches(4.2)

        for i, item in enumerate(content_items):
            x = left_x if i % 2 == 0 else right_x
            y = Inches(1.3 + (i // 2) * 1.1)

            # Icon circle
            shape = slide.shapes.add_shape(
                9,  # Oval
                x, y, Inches(0.4), Inches(0.4)
            )
            shape.fill.solid()
            shape.fill.fore_color.rgb = Colors.primary
            shape.line.fill.background()

            # Checkmark
            textbox = slide.shapes.add_textbox(x, y, Inches(0.4), Inches(0.4))
            text_frame = textbox.text_frame
            text_frame.text = "✓"
            p = text_frame.paragraphs[0]
            p.font.size = Pt(18)
            p.font.color.rgb = Colors.white
            p.alignment = PP_ALIGN.CENTER

            # Text
            textbox = slide.shapes.add_textbox(x + Inches(0.5), y, width - Inches(0.4), Inches(0.5))
            text_frame = textbox.text_frame
            text_frame.word_wrap = True
            p = text_frame.paragraphs[0]
            p.text = item
            p.font.size = Pt(14)
            p.font.color.rgb = Colors.dark
            p.space_before = Pt(0)
            p.space_after = Pt(0)
    else:
        # Single column with styled boxes
        for i, item in enumerate(content_items):
            y = Inches(1.3 + i * 1.3)

            # White box with shadow effect
            shape = slide.shapes.add_shape(
                1,  # Rectangle
                Inches(0.5), y, Inches(9), Inches(1.1)
            )
            shape.fill.solid()
            shape.fill.fore_color.rgb = Colors.white
            shape.line.color.rgb = RGBColor(200, 200, 200)

            # Text
            textbox = slide.shapes.add_textbox(Inches(0.7), y + Inches(0.05), Inches(8.6), Inches(1))
            text_frame = textbox.text_frame
            text_frame.word_wrap = True
            p = text_frame.paragraphs[0]
            p.text = item
            p.font.size = Pt(14)
            p.font.color.rgb = Colors.dark

def add_what_is_ai_slide(prs, blank_layout):
    """Slide 2: What is AI?"""
    slide = prs.slides.add_slide(blank_layout)

    # Light background
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = Colors.light

    # Accent bar
    shape = slide.shapes.add_shape(1, Inches(0), Inches(0), Inches(0.15), Inches(5.625))
    shape.fill.solid()
    shape.fill.fore_color.rgb = Colors.primary
    shape.line.fill.background()

    # Title
    textbox = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.7))
    text_frame = textbox.text_frame
    text_frame.text = "What is AI?"
    p = text_frame.paragraphs[0]
    p.font.size = Pt(44)
    p.font.color.rgb = Colors.dark
    p.font.bold = True

    # Definition box
    shape = slide.shapes.add_shape(1, Inches(0.5), Inches(1.3), Inches(9), Inches(2))
    shape.fill.solid()
    shape.fill.fore_color.rgb = Colors.white
    shape.line.color.rgb = Colors.primary
    shape.line.width = Pt(1)

    textbox = slide.shapes.add_textbox(Inches(0.8), Inches(1.5), Inches(8.4), Inches(1.6))
    text_frame = textbox.text_frame
    text_frame.word_wrap = True
    text_frame.text = "Artificial Intelligence is the simulation of human intelligence processes by machines, especially computer systems. It enables computers to learn, reason, perceive, and understand language."
    p = text_frame.paragraphs[0]
    p.font.size = Pt(18)
    p.font.color.rgb = Colors.dark
    p.alignment = PP_ALIGN.CENTER

    # Key characteristics (2x2 grid)
    keys = [
        ("Learning", "Improving from experience"),
        ("Reasoning", "Making logical decisions"),
        ("Perception", "Understanding inputs"),
        ("Language", "Processing communication")
    ]

    for i, (title, desc) in enumerate(keys):
        x = Inches(0.5 + i % 2 * 4.6)
        y = Inches(3.6 + i // 2 * 1.5)

        # Icon circle
        shape = slide.shapes.add_shape(9, x, y, Inches(0.5), Inches(0.5))
        shape.fill.solid()
        shape.fill.fore_color.rgb = Colors.primary
        shape.line.fill.background()

        textbox = slide.shapes.add_textbox(x, y, Inches(0.5), Inches(0.5))
        text_frame = textbox.text_frame
        text_frame.text = "✓"
        p = text_frame.paragraphs[0]
        p.font.size = Pt(20)
        p.font.color.rgb = Colors.white
        p.alignment = PP_ALIGN.CENTER

        # Title
        textbox = slide.shapes.add_textbox(x + Inches(0.6), y, Inches(4), Inches(0.3))
        text_frame = textbox.text_frame
        text_frame.text = title
        p = text_frame.paragraphs[0]
        p.font.size = Pt(14)
        p.font.color.rgb = Colors.dark
        p.font.bold = True

        # Description
        textbox = slide.shapes.add_textbox(x + Inches(0.6), y + Inches(0.35), Inches(4), Inches(0.3))
        text_frame = textbox.text_frame
        text_frame.text = desc
        p = text_frame.paragraphs[0]
        p.font.size = Pt(12)
        p.font.color.rgb = Colors.gray

def add_types_ai_slide(prs, blank_layout):
    """Slide 3: Types of AI"""
    slide = prs.slides.add_slide(blank_layout)

    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = Colors.light

    # Accent bar
    shape = slide.shapes.add_shape(1, Inches(0), Inches(0), Inches(0.15), Inches(5.625))
    shape.fill.solid()
    shape.fill.fore_color.rgb = Colors.primary
    shape.line.fill.background()

    # Title
    textbox = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.7))
    text_frame = textbox.text_frame
    text_frame.text = "Types of AI"
    p = text_frame.paragraphs[0]
    p.font.size = Pt(44)
    p.font.color.rgb = Colors.dark
    p.font.bold = True

    # Three columns
    types = [
        ("Narrow AI", "Designed for specific tasks", ["Voice assistants", "Recommendation systems", "Image recognition"], Colors.primary),
        ("Machine Learning", "Learns from data without explicit programming", ["Predictive analytics", "Spam filters", "Fraud detection"], Colors.secondary),
        ("Deep Learning", "Neural networks inspired by human brain", ["Language models", "Computer vision", "Game playing"], Colors.accent),
    ]

    card_width = Inches(2.8)
    card_gap = Inches(0.25)

    for i, (title, desc, examples, color) in enumerate(types):
        x = Inches(0.5 + i * (card_width + card_gap))
        y = Inches(1.3)

        # Card background
        shape = slide.shapes.add_shape(1, x, y, card_width, Inches(3.8))
        shape.fill.solid()
        shape.fill.fore_color.rgb = Colors.white
        shape.line.color.rgb = RGBColor(220, 220, 220)

        # Accent top bar
        shape = slide.shapes.add_shape(1, x, y, card_width, Inches(0.12))
        shape.fill.solid()
        shape.fill.fore_color.rgb = color
        shape.line.fill.background()

        # Title
        textbox = slide.shapes.add_textbox(x + Inches(0.2), y + Inches(0.25), card_width - Inches(0.4), Inches(0.4))
        text_frame = textbox.text_frame
        text_frame.text = title
        p = text_frame.paragraphs[0]
        p.font.size = Pt(18)
        p.font.color.rgb = Colors.dark
        p.font.bold = True

        # Description
        textbox = slide.shapes.add_textbox(x + Inches(0.2), y + Inches(0.65), card_width - Inches(0.4), Inches(0.6))
        text_frame = textbox.text_frame
        text_frame.word_wrap = True
        text_frame.text = desc
        p = text_frame.paragraphs[0]
        p.font.size = Pt(12)
        p.font.color.rgb = Colors.gray

        # Examples label
        textbox = slide.shapes.add_textbox(x + Inches(0.2), y + Inches(1.3), card_width - Inches(0.4), Inches(0.25))
        text_frame = textbox.text_frame
        text_frame.text = "Examples:"
        p = text_frame.paragraphs[0]
        p.font.size = Pt(11)
        p.font.color.rgb = color
        p.font.bold = True

        # Examples
        textbox = slide.shapes.add_textbox(x + Inches(0.2), y + Inches(1.55), card_width - Inches(0.4), Inches(2))
        text_frame = textbox.text_frame
        for j, example in enumerate(examples):
            if j > 0:
                p = text_frame.add_paragraph()
            else:
                p = text_frame.paragraphs[0]
            p.text = f"• {example}"
            p.font.size = Pt(11)
            p.font.color.rgb = Colors.dark
            p.space_before = Pt(2)
            p.space_after = Pt(2)

def add_how_ai_works_slide(prs, blank_layout):
    """Slide 4: How AI Works"""
    slide = prs.slides.add_slide(blank_layout)

    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = Colors.light

    # Accent bar
    shape = slide.shapes.add_shape(1, Inches(0), Inches(0), Inches(0.15), Inches(5.625))
    shape.fill.solid()
    shape.fill.fore_color.rgb = Colors.primary
    shape.line.fill.background()

    # Title
    textbox = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.7))
    text_frame = textbox.text_frame
    text_frame.text = "How AI Works"
    p = text_frame.paragraphs[0]
    p.font.size = Pt(44)
    p.font.color.rgb = Colors.dark
    p.font.bold = True

    # Process flow steps
    steps = [
        ("1", "Data Collection", "Gathering relevant information"),
        ("2", "Processing", "Cleaning and preparing data"),
        ("3", "Training", "Learning patterns from data"),
        ("4", "Inference", "Making predictions or decisions"),
    ]

    step_width = Inches(2.1)
    step_gap = Inches(0.25)

    for i, (num, title, desc) in enumerate(steps):
        x = Inches(0.5 + i * (step_width + step_gap))
        y = Inches(1.3)

        # Number circle
        shape = slide.shapes.add_shape(9, x, y, Inches(0.6), Inches(0.6))
        shape.fill.solid()
        shape.fill.fore_color.rgb = Colors.primary
        shape.line.fill.background()

        textbox = slide.shapes.add_textbox(x, y, Inches(0.6), Inches(0.6))
        text_frame = textbox.text_frame
        text_frame.text = num
        p = text_frame.paragraphs[0]
        p.font.size = Pt(28)
        p.font.color.rgb = Colors.white
        p.font.bold = True
        p.alignment = PP_ALIGN.CENTER

        # Title
        textbox = slide.shapes.add_textbox(x, y + Inches(0.8), step_width, Inches(0.4))
        text_frame = textbox.text_frame
        text_frame.text = title
        p = text_frame.paragraphs[0]
        p.font.size = Pt(14)
        p.font.color.rgb = Colors.dark
        p.font.bold = True
        p.alignment = PP_ALIGN.CENTER

        # Description
        textbox = slide.shapes.add_textbox(x, y + Inches(1.2), step_width, Inches(0.5))
        text_frame = textbox.text_frame
        text_frame.word_wrap = True
        text_frame.text = desc
        p = text_frame.paragraphs[0]
        p.font.size = Pt(11)
        p.font.color.rgb = Colors.gray
        p.alignment = PP_ALIGN.CENTER

    # Key insight box
    shape = slide.shapes.add_shape(1, Inches(0.5), Inches(3.5), Inches(9), Inches(1.6))
    shape.fill.solid()
    shape.fill.fore_color.rgb = Colors.primary
    shape.fill.fore_color.brightness = 0.2
    shape.line.color.rgb = Colors.primary
    shape.line.width = Pt(1)

    textbox = slide.shapes.add_textbox(Inches(0.8), Inches(3.7), Inches(8.4), Inches(1.2))
    text_frame = textbox.text_frame
    text_frame.word_wrap = True

    p = text_frame.paragraphs[0]
    p.text = "Key Insight:"
    p.font.size = Pt(14)
    p.font.color.rgb = Colors.dark
    p.font.bold = True

    p = text_frame.add_paragraph()
    p.text = "AI systems improve with more data and training time, enabling them to perform increasingly complex tasks with higher accuracy."
    p.font.size = Pt(12)
    p.font.color.rgb = Colors.dark

def add_applications_slide(prs, blank_layout):
    """Slide 5: AI Applications"""
    slide = prs.slides.add_slide(blank_layout)

    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = Colors.light

    # Accent bar
    shape = slide.shapes.add_shape(1, Inches(0), Inches(0), Inches(0.15), Inches(5.625))
    shape.fill.solid()
    shape.fill.fore_color.rgb = Colors.primary
    shape.line.fill.background()

    # Title
    textbox = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.7))
    text_frame = textbox.text_frame
    text_frame.text = "AI Applications"
    p = text_frame.paragraphs[0]
    p.font.size = Pt(44)
    p.font.color.rgb = Colors.dark
    p.font.bold = True

    # 2x3 grid of applications
    apps = [
        ("Healthcare", ["Disease diagnosis", "Drug discovery", "Personalized treatment"]),
        ("Finance", ["Fraud detection", "Algorithmic trading", "Risk assessment"]),
        ("Transportation", ["Autonomous vehicles", "Route optimization", "Traffic prediction"]),
        ("Entertainment", ["Content recommendation", "Game AI", "Creative generation"]),
        ("Education", ["Personalized learning", "Tutoring systems", "Automated grading"]),
        ("Business", ["Customer service", "Process automation", "Predictive analytics"]),
    ]

    app_width = Inches(2.85)
    app_height = Inches(1.9)
    x_gap = Inches(0.35)
    y_gap = Inches(0.3)

    for i, (title, items) in enumerate(apps):
        x = Inches(0.5 + (i % 2) * (app_width + x_gap))
        y = Inches(1.3 + (i // 2) * (app_height + y_gap))

        # Card background
        shape = slide.shapes.add_shape(1, x, y, app_width, app_height)
        shape.fill.solid()
        shape.fill.fore_color.rgb = Colors.white
        shape.line.color.rgb = RGBColor(220, 220, 220)

        # Icon circle
        shape = slide.shapes.add_shape(9, x + Inches(0.15), y + Inches(0.15), Inches(0.5), Inches(0.5))
        shape.fill.solid()
        shape.fill.fore_color.rgb = Colors.secondary
        shape.line.fill.background()

        # Use simple text icon since we can't do emoji reliably
        textbox = slide.shapes.add_textbox(x + Inches(0.15), y + Inches(0.15), Inches(0.5), Inches(0.5))
        text_frame = textbox.text_frame
        text_frame.text = "●"
        p = text_frame.paragraphs[0]
        p.font.size = Pt(24)
        p.font.color.rgb = Colors.white
        p.alignment = PP_ALIGN.CENTER

        # Title
        textbox = slide.shapes.add_textbox(x + Inches(0.8), y + Inches(0.15), app_width - Inches(1), Inches(0.3))
        text_frame = textbox.text_frame
        text_frame.text = title
        p = text_frame.paragraphs[0]
        p.font.size = Pt(14)
        p.font.color.rgb = Colors.dark
        p.font.bold = True

        # Items
        textbox = slide.shapes.add_textbox(x + Inches(0.2), y + Inches(0.75), app_width - Inches(0.4), Inches(1.1))
        text_frame = textbox.text_frame
        for j, item in enumerate(items):
            if j > 0:
                p = text_frame.add_paragraph()
            else:
                p = text_frame.paragraphs[0]
            p.text = f"• {item}"
            p.font.size = Pt(10)
            p.font.color.rgb = Colors.gray
            p.space_before = Pt(1)
            p.space_after = Pt(1)

def add_future_slide(prs, blank_layout):
    """Slide 6: The Future of AI"""
    slide = prs.slides.add_slide(blank_layout)

    # Dark background
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = Colors.dark

    # Decorative circle
    shape = slide.shapes.add_shape(9, Inches(-2), Inches(3.5), Inches(6), Inches(6))
    shape.fill.solid()
    shape.fill.fore_color.rgb = Colors.primary
    shape.fill.fore_color.brightness = 0.3
    shape.line.fill.background()

    # Title
    textbox = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.7))
    text_frame = textbox.text_frame
    text_frame.text = "The Future of AI"
    p = text_frame.paragraphs[0]
    p.font.size = Pt(44)
    p.font.color.rgb = Colors.white
    p.font.bold = True

    # Two sections
    sections = [
        ("Emerging Trends", ["Generative AI creating content", "AI-human collaboration", "Edge AI on devices", "Ethical AI development"], Colors.accent),
        ("Considerations", ["Privacy and data security", "Job market transformation", "Bias and fairness", "Regulation and governance"], Colors.secondary),
    ]

    section_width = Inches(4.25)
    section_gap = Inches(0.5)

    for i, (title, items, color) in enumerate(sections):
        x = Inches(0.5 + i * (section_width + section_gap))
        y = Inches(1.3)

        # Section box
        shape = slide.shapes.add_shape(1, x, y, section_width, Inches(3.2))
        shape.fill.solid()
        shape.fill.fore_color.rgb = Colors.white
        shape.fill.fore_color.brightness = -0.05
        shape.line.color.rgb = color
        shape.line.width = Pt(1)

        # Accent top bar
        shape = slide.shapes.add_shape(1, x, y, section_width, Inches(0.1))
        shape.fill.solid()
        shape.fill.fore_color.rgb = color
        shape.line.fill.background()

        # Title
        textbox = slide.shapes.add_textbox(x + Inches(0.2), y + Inches(0.2), section_width - Inches(0.4), Inches(0.35))
        text_frame = textbox.text_frame
        text_frame.text = title
        p = text_frame.paragraphs[0]
        p.font.size = Pt(16)
        p.font.color.rgb = Colors.white
        p.font.bold = True

        # Items
        textbox = slide.shapes.add_textbox(x + Inches(0.2), y + Inches(0.6), section_width - Inches(0.4), Inches(2.4))
        text_frame = textbox.text_frame
        for j, item in enumerate(items):
            if j > 0:
                p = text_frame.add_paragraph()
            else:
                p = text_frame.paragraphs[0]
            p.text = f"✓ {item}"
            p.font.size = Pt(11)
            p.font.color.rgb = RGBColor(226, 232, 240)
            p.space_before = Pt(2)
            p.space_after = Pt(2)

    # Closing statement
    textbox = slide.shapes.add_textbox(Inches(0.5), Inches(4.7), Inches(9), Inches(0.4))
    text_frame = textbox.text_frame
    text_frame.word_wrap = True
    text_frame.text = "AI will continue transforming how we live and work, creating new opportunities while requiring thoughtful consideration of its impact."
    p = text_frame.paragraphs[0]
    p.font.size = Pt(12)
    p.font.color.rgb = Colors.gray
    p.font.italic = True
    p.alignment = PP_ALIGN.CENTER

    # Thank you
    textbox = slide.shapes.add_textbox(Inches(0.5), Inches(5.15), Inches(9), Inches(0.4))
    text_frame = textbox.text_frame
    text_frame.text = "Thank You"
    p = text_frame.paragraphs[0]
    p.font.size = Pt(18)
    p.font.color.rgb = Colors.accent
    p.font.bold = True
    p.alignment = PP_ALIGN.CENTER

# Create all slides
add_title_slide(prs, blank_layout)
add_what_is_ai_slide(prs, blank_layout)
add_types_ai_slide(prs, blank_layout)
add_how_ai_works_slide(prs, blank_layout)
add_applications_slide(prs, blank_layout)
add_future_slide(prs, blank_layout)

# Save presentation
prs.save('AI-Introduction.pptx')
print("Presentation created: AI-Introduction.pptx")